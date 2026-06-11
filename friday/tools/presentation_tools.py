"""
Presentation generation tools — creates professional PowerPoint (.pptx) and
HTML slide decks from JSON, Markdown, or auto-generated demo content.

Libraries: python-pptx (optional, for PPTX output), webbrowser (stdlib).
"""
from __future__ import annotations

import datetime
import json
import os
import re
import textwrap
import webbrowser
from typing import Any, Optional

from friday.logging_utils import configure_logging
from friday._paths import FRIDAY_MEMORY

logger = configure_logging("presentation_tools")
PRESENTATION_DIR = os.path.join(FRIDAY_MEMORY, "presentations")
os.makedirs(PRESENTATION_DIR, exist_ok=True)

# ── Lazy dependency flags ─────────────────────────────────
HAS_PPTX = False
try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    HAS_PPTX = True
except ImportError:
    pass

# ── Theme definitions ─────────────────────────────────────

THEMES: dict[str, dict[str, Any]] = {
    "blue": {
        "bg": "#1a365d",
        "bg_gradient": "#2b6cb0",
        "text": "#ffffff",
        "accent": "#63b3ed",
        "heading_font": "Calibri Light",
        "body_font": "Calibri",
        "secondary_bg": "#ebf8ff",
        "slide_bg": "#ffffff",
        "slide_text": "#2d3748",
        "slide_accent": "#2b6cb0",
        "slide_heading": "#1a365d",
    },
    "dark": {
        "bg": "#0f0f0f",
        "bg_gradient": "#2d2d2d",
        "text": "#e0e0e0",
        "accent": "#bb86fc",
        "heading_font": "Segoe UI",
        "body_font": "Segoe UI",
        "secondary_bg": "#1e1e1e",
        "slide_bg": "#1e1e1e",
        "slide_text": "#e0e0e0",
        "slide_accent": "#bb86fc",
        "slide_heading": "#ffffff",
    },
    "green": {
        "bg": "#22543d",
        "bg_gradient": "#38a169",
        "text": "#ffffff",
        "accent": "#68d391",
        "heading_font": "Calibri Light",
        "body_font": "Calibri",
        "secondary_bg": "#f0fff4",
        "slide_bg": "#ffffff",
        "slide_text": "#2d3748",
        "slide_accent": "#38a169",
        "slide_heading": "#22543d",
    },
    "purple": {
        "bg": "#44337a",
        "bg_gradient": "#805ad5",
        "text": "#ffffff",
        "accent": "#b794f4",
        "heading_font": "Calibri Light",
        "body_font": "Calibri",
        "secondary_bg": "#faf5ff",
        "slide_bg": "#ffffff",
        "slide_text": "#2d3748",
        "slide_accent": "#805ad5",
        "slide_heading": "#44337a",
    },
    "corporate": {
        "bg": "#1e293b",
        "bg_gradient": "#334155",
        "text": "#ffffff",
        "accent": "#3b82f6",
        "heading_font": "Arial",
        "body_font": "Arial",
        "secondary_bg": "#f8fafc",
        "slide_bg": "#ffffff",
        "slide_text": "#1e293b",
        "slide_accent": "#3b82f6",
        "slide_heading": "#1e293b",
    },
}


# ── Helpers ────────────────────────────────────────────────

def _ts() -> str:
    return datetime.datetime.now().strftime("%Y%m%d_%H%M%S")


def _sanitize_filename(name: str) -> str:
    return re.sub(r"[^\w\- ]", "_", name).strip()[:60] or "presentation"


def _resolve_path(filename: str) -> str:
    return os.path.join(PRESENTATION_DIR, filename)


def _parse_markdown_slides(markdown_text: str) -> list[dict[str, Any]]:
    """Split markdown into slides by ## headings or --- separators."""
    slides: list[dict[str, Any]] = []
    raw_slides: list[str] = []

    # Try splitting by --- first
    parts = re.split(r"\n---+\n", markdown_text)
    for part in parts:
        part = part.strip()
        if not part:
            continue
        # Further split by ## headings within each --- block
        sub_parts = re.split(r"\n(?=##\s)", part)
        for sp in sub_parts:
            sp = sp.strip()
            if sp:
                raw_slides.append(sp)

    if not raw_slides:
        raw_slides = [markdown_text.strip()]

    for block in raw_slides:
        lines = block.split("\n")
        title = ""
        content_lines: list[str] = []
        for i, line in enumerate(lines):
            if line.startswith("## "):
                title = line[3:].strip()
            elif line.startswith("# ") and not title:
                title = line[2:].strip()
            else:
                content_lines.append(line)
        content = "\n".join(content_lines).strip()
        slide: dict[str, Any] = {"type": "content", "title": title or "Slide"}
        if content:
            slide["markdown_content"] = content
            # Parse markdown content for bullets, code blocks, images, tables
            bullets: list[str] = []
            code_blocks: list[str] = []
            images: list[str] = []
            tables: list[tuple[list[str], list[list[str]]]] = []
            current_code: list[str] = []
            in_code = False
            table_lines: list[str] = []
            in_table = False

            for cl in content.split("\n"):
                # Code block
                if cl.strip().startswith("```"):
                    if in_code:
                        code_blocks.append("\n".join(current_code))
                        current_code = []
                        in_code = False
                    else:
                        in_code = True
                    continue
                if in_code:
                    current_code.append(cl)
                    continue
                # Image
                img_match = re.match(r"!\[.*?\]\((.+?)\)", cl.strip())
                if img_match:
                    images.append(img_match.group(1))
                    continue
                # Table row
                if cl.strip().startswith("|"):
                    table_lines.append(cl.strip())
                    in_table = True
                    continue
                elif in_table and not cl.strip():
                    parsed = _parse_table_lines(table_lines)
                    if parsed:
                        tables.append(parsed)
                    table_lines = []
                    in_table = False
                    continue
                elif in_table:
                    table_lines.append(cl.strip())
                # Bullet
                if cl.strip().startswith("- ") or cl.strip().startswith("* "):
                    bullets.append(cl.strip()[2:])
                elif cl.strip().startswith("1. "):
                    bullets.append(cl.strip())
                elif cl.strip() and not cl.strip().startswith("!"):
                    bullets.append(cl.strip())

            if in_table and table_lines:
                parsed = _parse_table_lines(table_lines)
                if parsed:
                    tables.append(parsed)

            if bullets:
                slide["bullets"] = bullets
            if code_blocks:
                slide["code"] = code_blocks
            if images:
                slide["images"] = images
            if tables:
                slide["tables"] = tables

        slides.append(slide)

    return slides


def _parse_table_lines(lines: list[str]) -> tuple[list[str], list[list[str]]]:
    """Parse markdown table lines into (headers, rows)."""
    if len(lines) < 2:
        return [], []
    headers = [h.strip().strip("|") for h in lines[0].split("|") if h.strip()]
    rows: list[list[str]] = []
    for line in lines[2:]:
        cells = [c.strip().strip("|") for c in line.split("|") if c.strip()]
        if cells:
            rows.append(cells)
    return headers, rows


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


def _add_pptx_text_box(slide, left: int, top: int, width: int, height: int,
                       text: str, font_name: str = "Calibri", font_size: int = 18,
                       bold: bool = False, color: str = "#2d3748",
                       alignment: str = "left") -> None:
    """Add a text box to a PPTX slide with formatting."""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*_hex_to_rgb(color))
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
    tf.paragraphs[0].space_after = Pt(4)
    return txBox


def _add_pptx_bullets(slide, left: int, top: int, width: int, height: int,
                      items: list[str], font_name: str = "Calibri",
                      font_size: int = 16, color: str = "#2d3748",
                      accent_color: str = "#2b6cb0") -> None:
    """Add a bulleted list to a PPTX slide."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(item)
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*_hex_to_rgb(color))
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        # Add bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "–"})
        # Remove existing buNone if any
        for existing_buNone in pPr.findall(qn("a:buNone")):
            pPr.remove(existing_buNone)
        pPr.append(buChar)
    return txBox


def _build_pptx(slides_data: list[dict[str, Any]], title: str,
                theme_name: str = "blue") -> Optional[str]:
    """Build a PPTX file from slide data. Returns file path or None."""
    if not HAS_PPTX:
        return None

    theme = THEMES.get(theme_name, THEMES["blue"])
    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, slide_data in enumerate(slides_data):
        st = slide_data.get("type", "content")
        notes_text = slide_data.get("notes", "")

        if st == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["bg"]))
            # Title
            _add_pptx_text_box(slide, 1, 2.0, 11.333, 2.0,
                               slide_data.get("title", title),
                               font_name=theme["heading_font"], font_size=44,
                               bold=True, color=theme["text"], alignment="center")
            # Subtitle
            sub = slide_data.get("subtitle", "")
            if sub:
                _add_pptx_text_box(slide, 1, 4.2, 11.333, 1.5, sub,
                                   font_name=theme["body_font"], font_size=22,
                                   bold=False, color=theme["accent"], alignment="center")
            # Accent line
            from pptx.util import Emu
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.5), Inches(3.9), Inches(4.333), Inches(0.04)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["accent"]))
            shape.line.fill.background()

        elif st == "content":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_bg"]))
            # Top accent bar
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(13.333), Inches(0.06)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_accent"]))
            bar.line.fill.background()
            # Title
            _add_pptx_text_box(slide, 0.8, 0.4, 11.733, 0.8,
                               slide_data.get("title", ""),
                               font_name=theme["heading_font"], font_size=30,
                               bold=True, color=theme["slide_heading"], alignment="left")
            # Divider line under title
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.03)
            )
            div.fill.solid()
            div.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_accent"]))
            div.line.fill.background()
            # Bullets
            bullets = slide_data.get("bullets") or []
            if bullets:
                _add_pptx_bullets(slide, 0.8, 1.6, 11.733, 5.0,
                                  bullets,
                                  font_name=theme["body_font"], font_size=18,
                                  color=theme["slide_text"],
                                  accent_color=theme["slide_accent"])

        elif st == "two_column":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_bg"]))
            # Title
            _add_pptx_text_box(slide, 0.8, 0.3, 11.733, 0.8,
                               slide_data.get("title", ""),
                               font_name=theme["heading_font"], font_size=28,
                               bold=True, color=theme["slide_heading"])
            # Left column header
            _add_pptx_text_box(slide, 0.8, 1.3, 5.5, 0.5,
                               slide_data.get("left_header", "Left"),
                               font_name=theme["heading_font"], font_size=20,
                               bold=True, color=theme["slide_accent"])
            # Left bullets
            left_items = slide_data.get("left", [])
            if left_items:
                _add_pptx_bullets(slide, 0.8, 1.9, 5.5, 4.5,
                                  left_items,
                                  font_name=theme["body_font"], font_size=16,
                                  color=theme["slide_text"])
            # Right column header
            _add_pptx_text_box(slide, 7.0, 1.3, 5.5, 0.5,
                               slide_data.get("right_header", "Right"),
                               font_name=theme["heading_font"], font_size=20,
                               bold=True, color=theme["slide_accent"])
            # Right bullets
            right_items = slide_data.get("right", [])
            if right_items:
                _add_pptx_bullets(slide, 7.0, 1.9, 5.5, 4.5,
                                  right_items,
                                  font_name=theme["body_font"], font_size=16,
                                  color=theme["slide_text"])
            # Vertical divider
            vdiv = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(6.55), Inches(1.5), Inches(0.02), Inches(5.0)
            )
            vdiv.fill.solid()
            vdiv.fill.fore_color.rgb = RGBColor(*_hex_to_rgb("#e2e8f0"))
            vdiv.line.fill.background()

        elif st == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_bg"]))
            # Title
            _add_pptx_text_box(slide, 0.8, 0.3, 11.733, 0.8,
                               slide_data.get("title", ""),
                               font_name=theme["heading_font"], font_size=28,
                               bold=True, color=theme["slide_heading"])
            # Image
            img_url = slide_data.get("image_url", "")
            if img_url and os.path.exists(img_url):
                try:
                    slide.shapes.add_picture(img_url,
                                             Inches(1.5), Inches(1.5),
                                             width=Inches(10), height=Inches(5))
                except Exception:
                    pass
            # Caption
            caption = slide_data.get("caption", "")
            if caption:
                _add_pptx_text_box(slide, 1.5, 6.6, 10, 0.6,
                                   caption,
                                   font_name=theme["body_font"], font_size=14,
                                   bold=False, color=theme["slide_accent"],
                                   alignment="center")

        elif st == "table":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_bg"]))
            # Title
            _add_pptx_text_box(slide, 0.8, 0.3, 11.733, 0.8,
                               slide_data.get("title", ""),
                               font_name=theme["heading_font"], font_size=28,
                               bold=True, color=theme["slide_heading"])
            # Table
            headers = slide_data.get("headers", [])
            rows = slide_data.get("rows", [])
            if headers and rows:
                num_cols = len(headers)
                num_rows = len(rows) + 1
                table_shape = slide.shapes.add_table(num_rows, num_cols,
                                                     Inches(0.8), Inches(1.5),
                                                     Inches(11.733), Inches(5.5))
                table = table_shape.table
                # Header row
                for ci, h in enumerate(headers):
                    cell = table.cell(0, ci)
                    cell.text = str(h)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(*_hex_to_rgb(theme["slide_bg"]))
                        paragraph.font.name = theme["heading_font"]
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme["slide_accent"]))
                # Data rows
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        cell = table.cell(ri + 1, ci)
                        cell.text = str(val)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.font.color.rgb = RGBColor(*_hex_to_rgb(theme["slide_text"]))
                            paragraph.font.name = theme["body_font"]
                        if ri % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(*_hex_to_rgb("#f7fafc"))

        # Add notes
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
            except Exception:
                pass

    filename = f"{_sanitize_filename(title)}_{_ts()}.pptx"
    filepath = _resolve_path(filename)
    prs.save(filepath)
    logger.info("Created PPTX: %s (%d slides)", filepath, len(slides_data))
    return filepath


# ── HTML slide deck generator ─────────────────────────────

_HTML_TEMPLATE = """\
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ font-family: '{font_family}', sans-serif; overflow: hidden;
         background: {bg_color}; height: 100vh; }}
  .deck {{ position: relative; width: 100vw; height: 100vh; overflow: hidden; }}
  .slides {{ position: relative; width: 100%; height: 100%; }}
  .slide {{ position: absolute; top: 0; left: 0; width: 100%; height: 100%;
           display: flex; flex-direction: column; justify-content: center;
           align-items: center; padding: 60px 80px;
           opacity: 0; transform: scale(0.95);
           transition: opacity 0.5s ease, transform 0.5s ease;
           background: {slide_bg}; color: {slide_text}; }}
  .slide.active {{ opacity: 1; transform: scale(1); z-index: 10; }}
  .slide-title {{ font-size: 2.8em; font-weight: 700; margin-bottom: 20px;
                  color: {slide_heading}; text-align: center; }}
  .slide-subtitle {{ font-size: 1.3em; color: {accent_color}; margin-bottom: 30px;
                     text-align: center; }}
  .slide-content {{ font-size: 1.2em; line-height: 1.6; max-width: 900px;
                    width: 100%; text-align: left; }}
  .slide-content ul {{ list-style: none; padding: 0; }}
  .slide-content ul li {{ padding: 8px 0 8px 30px; position: relative;
                          font-size: 1.15em; }}
  .slide-content ul li::before {{ content: "–"; color: {accent_color};
                                  font-weight: bold;
                                  position: absolute; left: 8px; }}
  .slide-content code {{ background: {code_bg}; padding: 2px 6px; border-radius: 4px;
                         font-family: 'Courier New', monospace; font-size: 0.85em; }}
  pre {{ background: {code_bg}; padding: 16px 20px; border-radius: 8px;
         overflow-x: auto; font-size: 0.85em; line-height: 1.5;
         border-left: 3px solid {accent_color}; margin: 12px 0; }}
  pre code {{ background: none; padding: 0; }}
  table {{ width: 100%; border-collapse: collapse; margin: 16px 0;
           font-size: 0.9em; }}
  th {{ background: {accent_color}; color: #fff; padding: 10px 14px;
        text-align: left; font-weight: 600; }}
  td {{ padding: 8px 14px; border-bottom: 1px solid {border_color}; }}
  tr:nth-child(even) td {{ background: {alt_row}; }}
  img {{ max-width: 80%; max-height: 50vh; border-radius: 8px;
         margin: 16px 0; box-shadow: 0 4px 20px rgba(0,0,0,0.1); }}
  .caption {{ font-size: 0.85em; color: {accent_color}; text-align: center;
              margin-top: 4px; }}
  .nav {{ position: fixed; bottom: 30px; left: 50%; transform: translateX(-50%);
          display: flex; align-items: center; gap: 20px; z-index: 100;
          background: {nav_bg}; padding: 10px 24px; border-radius: 30px;
          box-shadow: 0 2px 20px rgba(0,0,0,0.15); }}
  .nav button {{ background: {accent_color}; color: #fff; border: none;
                 width: 40px; height: 40px; border-radius: 50%;
                 font-size: 1.3em; cursor: pointer;
                 transition: background 0.2s, transform 0.2s; }}
  .nav button:hover {{ background: {accent_hover}; transform: scale(1.1); }}
  .nav button:disabled {{ opacity: 0.4; cursor: default; transform: none; }}
  .nav .counter {{ font-size: 0.95em; color: {counter_color};
                   min-width: 100px; text-align: center; }}
  .title-slide {{ background: {title_bg}; }}
  .title-slide .slide-title {{ font-size: 3.5em; color: {title_text}; }}
  .title-slide .slide-subtitle {{ color: {title_accent}; }}
  @media (max-width: 768px) {{
    .slide {{ padding: 30px 24px; }}
    .slide-title {{ font-size: 1.8em; }}
    .slide-content {{ font-size: 1em; }}
    .title-slide .slide-title {{ font-size: 2.2em; }}
  }}
</style>
</head>
<body>
<div class="deck">
  <div class="slides" id="slides">
    {slides_html}
  </div>
  <div class="nav">
    <button id="prevBtn" onclick="navigate(-1)" aria-label="Previous slide">◀</button>
    <span class="counter" id="slideCounter">Slide 1 of {total}</span>
    <button id="nextBtn" onclick="navigate(1)" aria-label="Next slide">▶</button>
  </div>
</div>
<script>
  let current = 0;
  const slides = document.querySelectorAll('.slide');
  const total = slides.length;
  const counter = document.getElementById('slideCounter');
  const prevBtn = document.getElementById('prevBtn');
  const nextBtn = document.getElementById('nextBtn');

  function showSlide(idx) {{
    slides.forEach((s, i) => {{
      s.classList.toggle('active', i === idx);
    }});
    counter.textContent = `Slide ${{idx + 1}} of ${{total}}`;
    prevBtn.disabled = idx === 0;
    nextBtn.disabled = idx === total - 1;
  }}

  function navigate(dir) {{
    const next = Math.max(0, Math.min(total - 1, current + dir));
    if (next !== current) {{ current = next; showSlide(current); }}
  }}

  document.addEventListener('keydown', (e) => {{
    if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') navigate(-1);
    if (e.key === 'ArrowRight' || e.key === 'ArrowDown' || e.key === ' ') {{ e.preventDefault(); navigate(1); }}
    if (e.key === 'Home') {{ current = 0; showSlide(0); }}
    if (e.key === 'End') {{ current = total - 1; showSlide(total - 1); }}
  }});

  // Touch support
  let touchStartX = 0;
  document.addEventListener('touchstart', (e) => {{ touchStartX = e.changedTouches[0].screenX; }});
  document.addEventListener('touchend', (e) => {{
    const diff = touchStartX - e.changedTouches[0].screenX;
    if (Math.abs(diff) > 50) navigate(diff > 0 ? 1 : -1);
  }});

  showSlide(0);
</script>
</body>
</html>"""


def _build_html_slide(slide_data: dict[str, Any], theme: dict[str, Any],
                      idx: int) -> str:
    """Build a single HTML slide element."""
    st = slide_data.get("type", "content")
    is_title = st == "title"
    extra_class = " title-slide" if is_title else ""
    title_text = slide_data.get("title", "")
    subtitle_text = slide_data.get("subtitle", "")
    notes_text = slide_data.get("notes", "")

    parts = [f'<div class="slide{extra_class}" data-index="{idx}">']

    if is_title:
        parts.append(f'<div class="slide-title">{title_text}</div>')
        if subtitle_text:
            parts.append(f'<div class="slide-subtitle">{subtitle_text}</div>')
    else:
        parts.append(f'<div class="slide-title">{title_text}</div>')
        parts.append('<div class="slide-content">')

        bullets = slide_data.get("bullets") or []
        code_blocks = slide_data.get("code") or []
        images = slide_data.get("images") or []
        tables_data = slide_data.get("tables") or []
        caption = slide_data.get("caption", "")
        left_items = slide_data.get("left") or []
        right_items = slide_data.get("right") or []
        headers = slide_data.get("headers") or []
        rows = slide_data.get("rows") or []
        md_content = slide_data.get("markdown_content", "")

        # Two column layout
        if st == "two_column":
            parts.append('<div style="display:flex;gap:40px;width:100%;">')
            left_h = slide_data.get("left_header", "Left")
            right_h = slide_data.get("right_header", "Right")
            parts.append('<div style="flex:1;">')
            parts.append(f'<h3 style="color:{theme["slide_accent"]};margin-bottom:12px;">{left_h}</h3>')
            parts.append('<ul>')
            for item in left_items:
                parts.append(f'<li>{item}</li>')
            parts.append('</ul></div>')
            parts.append('<div style="flex:1;">')
            parts.append(f'<h3 style="color:{theme["slide_accent"]};margin-bottom:12px;">{right_h}</h3>')
            parts.append('<ul>')
            for item in right_items:
                parts.append(f'<li>{item}</li>')
            parts.append('</ul></div>')
            parts.append('</div>')

        # Table
        elif st == "table" and headers and rows:
            parts.append('<table><thead><tr>')
            for h in headers:
                parts.append(f'<th>{h}</th>')
            parts.append('</tr></thead><tbody>')
            for row in rows:
                parts.append('<tr>')
                for cell in row:
                    parts.append(f'<td>{cell}</td>')
                parts.append('</tr>')
            parts.append('</tbody></table>')

        # Image
        elif st == "image":
            img_url = slide_data.get("image_url", "")
            if img_url:
                parts.append(f'<img src="{img_url}" alt="{caption or title_text}" />')
            if caption:
                parts.append(f'<div class="caption">{caption}</div>')

        # Content with bullets, code, images, tables from markdown
        else:
            # Show bullets
            if bullets:
                parts.append('<ul>')
                for b in bullets:
                    # Check for image in bullet
                    img_match = re.match(r"!\[.*?\]\((.+?)\)", b)
                    if img_match:
                        parts.append(f'<img src="{img_match.group(1)}" style="max-width:100%;" />')
                    else:
                        parts.append(f'<li>{b}</li>')
                parts.append('</ul>')

            # Show code blocks
            for code_text in code_blocks:
                escaped = (code_text
                           .replace("&", "&amp;")
                           .replace("<", "&lt;")
                           .replace(">", "&gt;"))
                parts.append(f'<pre><code>{escaped}</code></pre>')

            # Show images
            for img in images:
                parts.append(f'<img src="{img}" alt="Slide image" />')

            # Show tables
            for hdrs, tbl_rows in tables_data:
                parts.append('<table><thead><tr>')
                for h in hdrs:
                    parts.append(f'<th>{h}</th>')
                parts.append('</tr></thead><tbody>')
                for row in tbl_rows:
                    parts.append('<tr>')
                    for cell in row:
                        parts.append(f'<td>{cell}</td>')
                    parts.append('</tr>')
                parts.append('</tbody></table>')

        parts.append('</div>')

        # Notes as hidden data attribute
        if notes_text:
            safe_notes = (notes_text
                          .replace("&", "&amp;")
                          .replace("<", "&lt;")
                          .replace(">", "&gt;")
                          .replace('"', "&quot;"))
            parts.append(f'<div style="display:none" class="notes">{safe_notes}</div>')

    parts.append('</div>')
    return "\n".join(parts)


def _build_html(slides_data: list[dict[str, Any]], title: str,
                theme_name: str = "blue") -> Optional[str]:
    """Build an HTML slide deck file. Returns file path or None."""
    theme = THEMES.get(theme_name, THEMES["blue"])
    slides_html_parts: list[str] = []
    for idx, sd in enumerate(slides_data):
        slides_html_parts.append(_build_html_slide(sd, theme, idx))

    slides_html = "\n".join(slides_html_parts)

    html = _HTML_TEMPLATE.format(
        title=title,
        font_family=theme.get("body_font", "Segoe UI"),
        bg_color=theme["bg"],
        slide_bg=theme["slide_bg"],
        slide_text=theme["slide_text"],
        slide_heading=theme["slide_heading"],
        accent_color=theme["slide_accent"],
        code_bg="#1e293b" if theme_name == "dark" else "#f1f5f9",
        border_color="#e2e8f0",
        alt_row="#f8fafc",
        nav_bg=theme["bg"],
        counter_color=theme["text"],
        accent_hover=theme.get("accent", theme["slide_accent"]),
        title_bg=theme["bg"],
        title_text=theme["text"],
        title_accent=theme["accent"],
        slides_html=slides_html,
        total=len(slides_data),
    )

    filename = f"{_sanitize_filename(title)}_{_ts()}.html"
    filepath = _resolve_path(filename)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(html)
    logger.info("Created HTML deck: %s (%d slides)", filepath, len(slides_data))
    return filepath


# ── Public API ─────────────────────────────────────────────

def presentation_create(slides_json: str, title: str = "Presentation",
                        format: str = "pptx") -> dict[str, Any]:
    """Create a presentation from a JSON description of slides.

    Args:
        slides_json: JSON string array of slide objects.
        title: Presentation title.
        format: "pptx" or "html".

    Returns:
        dict with success, file_path, format, slide_count, title.
    """
    try:
        slides_data = json.loads(slides_json) if isinstance(slides_json, str) else slides_json
    except json.JSONDecodeError as e:
        return {"success": False, "error": f"Invalid JSON: {e}"}

    if not isinstance(slides_data, list):
        return {"success": False, "error": "slides_json must be a JSON array"}

    # Detect theme from first slide
    theme_name = slides_data[0].get("theme", "blue") if slides_data else "blue"
    if theme_name not in THEMES:
        theme_name = "blue"

    slide_count = len(slides_data)

    if format == "pptx":
        if not HAS_PPTX:
            return {
                "success": False,
                "error": "python-pptx not installed. Install with: pip install python-pptx",
                "install_hint": "pip install python-pptx",
            }
        filepath = _build_pptx(slides_data, title, theme_name)
        if not filepath:
            return {"success": False, "error": "Failed to create PPTX"}
    elif format == "html":
        filepath = _build_html(slides_data, title, theme_name)
        if not filepath:
            return {"success": False, "error": "Failed to create HTML"}
    else:
        return {"success": False, "error": f"Unknown format: {format}. Use 'pptx' or 'html'."}

    try:
        webbrowser.open(filepath)
    except Exception:
        pass

    return {
        "success": True,
        "file_path": filepath,
        "format": format,
        "slide_count": slide_count,
        "title": title,
    }


def presentation_create_from_markdown(markdown_text: str,
                                      title: str = "Presentation",
                                      format: str = "html") -> dict[str, Any]:
    """Parse markdown into slides and generate a presentation.

    Splits by ## headings or --- separators. Supports bullets,
    code blocks, images, and tables within each slide.

    Args:
        markdown_text: Raw markdown string.
        title: Presentation title.
        format: "pptx" or "html".

    Returns:
        dict with success, file_path, format, slide_count, title.
    """
    if not markdown_text or not markdown_text.strip():
        return {"success": False, "error": "Empty markdown text"}

    slides_data = _parse_markdown_slides(markdown_text)

    if not slides_data:
        return {"success": False, "error": "No slide content found in markdown"}

    # Prepend a title slide
    slides_data.insert(0, {"type": "title", "title": title, "subtitle": "Generated from Markdown"})

    return presentation_create(json.dumps(slides_data), title=title, format=format)


def presentation_create_demo(topic: str = "FRIDAY AI Assistant",
                             slides: int = 8) -> dict[str, Any]:
    """Create a sample/demo presentation about the given topic.

    Generates N slides with a template structure: title, agenda,
    content slides, summary, thank you.

    Args:
        topic: Topic for the demo presentation.
        slides: Number of slides (minimum 4).

    Returns:
        dict with success, file_path, format, slide_count, title.
    """
    n = max(4, min(slides, 30))

    demo_slides: list[dict[str, Any]] = [
        {
            "type": "title",
            "title": topic,
            "subtitle": "Presented by FRIDAY · Stark Industries AI",
        },
        {
            "type": "content",
            "title": "Agenda",
            "bullets": [
                "Overview of " + topic,
                "Key Features & Capabilities",
                "Architecture & Design",
                "Use Cases & Applications",
                "Performance & Metrics",
                "Roadmap & Future Work",
                "Q & A",
            ],
            "notes": "Welcome everyone. Today we'll cover these key areas.",
        },
    ]

    content_ideas = [
        ("What is " + topic + "?",
         [
             "Next-generation intelligent system",
             "Built on cutting-edge AI technology",
             "Designed for scalability and reliability",
             "Seamless integration with existing workflows",
         ]),
        ("Core Capabilities",
         [
             "Natural language understanding & generation",
             "Multi-modal processing (text, image, audio)",
             "Real-time data analysis & visualization",
             "Autonomous task execution & orchestration",
             "Context-aware memory & learning",
         ]),
        ("Architecture Overview",
         [
             "Modular microservices architecture",
             "Event-driven communication bus",
             "Distributed agent system with heartbeats",
             "Vector memory & knowledge graph storage",
             "Secure sandboxed execution environment",
         ]),
        ("Key Integrations",
         [
             "Web browsing & automation (browser-use)",
             "Desktop application control",
             "Email & calendar management",
             "Cloud services (Google, AWS, Azure)",
             "IoT & smart home devices",
             "Code execution & version control",
         ]),
        ("Performance Highlights",
         [
             "Sub-100ms response time for common queries",
             "99.9% uptime with auto-recovery",
             "Handles 1000+ concurrent sessions",
             "Multi-model AI routing for optimal results",
             "End-to-end encryption for all data",
         ]),
        ("Use Cases",
         [
             "Personal AI assistant for daily productivity",
             "Enterprise automation & workflow management",
             "Research & data analysis platform",
             "Cybersecurity threat detection & response",
             "Content generation & media production",
         ]),
        ("Future Roadmap",
         [
             "Enhanced multi-agent collaboration",
             "Improved long-term memory & reasoning",
             "Expanded third-party integrations",
             "Mobile & edge device deployment",
             "Advanced vision & multimodal capabilities",
         ]),
    ]

    # Add content slides (fit to requested count)
    for i in range(min(len(content_ideas), n - 3)):
        ci = content_ideas[i]
        demo_slides.append({
            "type": "content",
            "title": ci[0],
            "bullets": ci[1],
            "notes": f"Slide discussing: {ci[0]}",
        })

    # Adjust to exact count
    while len(demo_slides) < n - 2:
        idx = len(demo_slides) - 2
        ci = content_ideas[idx % len(content_ideas)]
        demo_slides.append({
            "type": "content",
            "title": ci[0] + " (cont.)",
            "bullets": ci[1],
            "notes": "",
        })

    # Summary
    demo_slides.append({
        "type": "content",
        "title": "Key Takeaways",
        "bullets": [
            f"{topic} represents a major leap in AI assistance",
            "Modular, secure, and highly extensible architecture",
            "Proven performance across diverse use cases",
            "Continuous improvement driven by community feedback",
        ],
        "notes": "In summary, our platform delivers unmatched capabilities.",
    })

    # Thank you
    demo_slides.append({
        "type": "title",
        "title": "Thank You",
        "subtitle": f"Questions? · {topic} · FRIDAY by Stark Industries",
    })

    return presentation_create(json.dumps(demo_slides), title=topic, format="html")


def presentation_list(limit: int = 20) -> dict[str, Any]:
    """List all generated presentations.

    Args:
        limit: Maximum number of entries.

    Returns:
        dict with success, presentations (list of {id, name, path, size, modified}).
    """
    if not os.path.isdir(PRESENTATION_DIR):
        return {"success": True, "presentations": []}

    entries: list[dict[str, Any]] = []
    try:
        for fname in sorted(os.listdir(PRESENTATION_DIR), reverse=True):
            fpath = os.path.join(PRESENTATION_DIR, fname)
            if not os.path.isfile(fpath):
                continue
            ext = os.path.splitext(fname)[1].lower()
            if ext not in (".pptx", ".html", ".htm"):
                continue
            stat = os.stat(fpath)
            entries.append({
                "id": fname,
                "name": os.path.splitext(fname)[0],
                "path": fpath,
                "size": stat.st_size,
                "modified": datetime.datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "format": "pptx" if ext == ".pptx" else "html",
            })
    except OSError as e:
        return {"success": False, "error": str(e)}

    entries = entries[:limit]
    return {"success": True, "presentations": entries, "count": len(entries)}


def presentation_open(presentation_id: str) -> dict[str, Any]:
    """Open a specific presentation file.

    Args:
        presentation_id: Filename (e.g. "MyPres_20250101_120000.pptx").

    Returns:
        dict with success, file_path.
    """
    sanitized = os.path.basename(presentation_id)
    fpath = _resolve_path(sanitized)
    if not os.path.isfile(fpath):
        return {"success": False, "error": f"Presentation not found: {sanitized}",
                "search_hint": f"Check {PRESENTATION_DIR}/"}

    try:
        webbrowser.open(fpath)
    except Exception as e:
        return {"success": False, "error": f"Failed to open file: {e}"}

    return {"success": True, "file_path": fpath}


def presentation_export_to_html(pptx_path: str) -> dict[str, Any]:
    """Export a PPTX file to an HTML slide deck.

    Extracts text from each slide using python-pptx and creates
    an HTML presentation with navigation.

    Args:
        pptx_path: Path to the .pptx file.

    Returns:
        dict with success, file_path, slide_count.
    """
    if not HAS_PPTX:
        return {
            "success": False,
            "error": "python-pptx not installed. Install with: pip install python-pptx",
            "install_hint": "pip install python-pptx",
        }
    if not os.path.isfile(pptx_path):
        return {"success": False, "error": f"File not found: {pptx_path}"}

    try:
        prs = PptxPresentation(pptx_path)
    except Exception as e:
        return {"success": False, "error": f"Failed to read PPTX: {e}"}

    slides_data: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        slide_dict: dict[str, Any] = {"type": "content", "title": f"Slide {i + 1}"}
        text_frames: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_frames.append(t)
            if shape.has_table:
                table = shape.table
                headers: list[str] = []
                rows: list[list[str]] = []
                for ri, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    if ri == 0:
                        headers = cells
                    else:
                        rows.append(cells)
                if headers:
                    slide_dict["headers"] = headers
                if rows:
                    slide_dict["rows"] = rows

        if text_frames:
            # First text is usually the title
            slide_dict["title"] = text_frames[0]
            if len(text_frames) > 1:
                slide_dict["bullets"] = text_frames[1:]
        else:
            slide_dict["title"] = f"Slide {i + 1}"

        slides_data.append(slide_dict)

    title = os.path.splitext(os.path.basename(pptx_path))[0]
    filepath = _build_html(slides_data, title)
    if not filepath:
        return {"success": False, "error": "Failed to generate HTML"}

    try:
        webbrowser.open(filepath)
    except Exception:
        pass

    return {
        "success": True,
        "file_path": filepath,
        "format": "html",
        "slide_count": len(slides_data),
        "title": title,
    }
