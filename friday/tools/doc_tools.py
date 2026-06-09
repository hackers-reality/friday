"""
Document Processing tools
Libraries: python-docx, openpyxl, pandas, polars, pdfplumber, pypdf, reportlab, python-pptx
"""
import asyncio
import os
import tempfile
from typing import Any

# ── Word (python-docx) ──
HAS_DOCX = False
try:
    from docx import Document
    from docx.shared import Inches, Pt
    HAS_DOCX = True
except ImportError:
    pass


async def read_docx(path: str) -> dict[str, Any]:
    if not HAS_DOCX:
        return {"error": "python-docx not installed"}
    if not os.path.exists(path):
        return {"error": f"File not found: {path}"}
    try:
        doc = await asyncio.get_event_loop().run_in_executor(None, lambda: Document(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return {"path": path, "paragraphs": len(paragraphs), "text": "\n".join(paragraphs[:500]),
                "tables": len(doc.tables), "sections": len(doc.sections)}
    except Exception as e:
        return {"error": str(e)}


async def create_docx(content: str | list[dict[str, Any]] | None = None,
                       sections: list[dict[str, Any]] | None = None,
                       output_path: str | None = None) -> dict[str, Any]:
    """Create a Word document with optional sections (headings, paragraphs, tables, charts)."""
    if not HAS_DOCX:
        return {"error": "python-docx not installed"}
    try:
        doc = Document()
        secs = sections or ([{"type": "paragraph", "text": content}] if isinstance(content, str) else [])
        for sec in secs:
            t = sec.get("type", "paragraph")
            if t == "heading":
                doc.add_heading(sec.get("text", ""), level=sec.get("level", 1))
            elif t == "paragraph":
                p = doc.add_paragraph(sec.get("text", ""))
                if sec.get("bold"):
                    for run in p.runs:
                        run.bold = True
            elif t == "bullet_list":
                for item in sec.get("items", []):
                    doc.add_paragraph(item, style="List Bullet")
            elif t == "numbered_list":
                for item in sec.get("items", []):
                    doc.add_paragraph(item, style="List Number")
            elif t == "table":
                rows = sec.get("rows", [])
                if rows:
                    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                    for i, row in enumerate(rows):
                        for j, cell in enumerate(row):
                            table.cell(i, j).text = str(cell)
                    doc.add_paragraph()
            elif t == "chart":
                chart_path = _make_chart(
                    sec.get("chart_type", "bar"), sec.get("data", []),
                    title=sec.get("title", ""),
                    xlabel=sec.get("xlabel", ""),
                    ylabel=sec.get("ylabel", ""),
                    labels=sec.get("labels"),
                    data2=sec.get("data2"),
                    data3=sec.get("data3"),
                )
                if chart_path:
                    doc.add_picture(chart_path, width=Inches(5.5))
                    doc.add_paragraph()
            elif t == "divider":
                doc.add_paragraph("─" * 50)
            elif t == "code":
                from docx.shared import RGBColor
                p = doc.add_paragraph()
                run = p.add_run(sec.get("text", ""))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)
        out = output_path or os.path.join(tempfile.gettempdir(), "friday_document.docx")
        await asyncio.get_event_loop().run_in_executor(None, lambda: doc.save(out))
        return {"path": out, "sections": len(secs)}
    except Exception as e:
        return {"error": str(e)}


# ── Excel (openpyxl) ──
HAS_OPENPYXL = False
try:
    from openpyxl import load_workbook, Workbook
    HAS_OPENPYXL = True
except ImportError:
    pass


async def read_excel(path: str, sheet: str | None = None) -> dict[str, Any]:
    if not HAS_OPENPYXL:
        return {"error": "openpyxl not installed"}
    if not os.path.exists(path):
        return {"error": f"File not found: {path}"}
    try:
        wb = await asyncio.get_event_loop().run_in_executor(None, lambda: load_workbook(path, data_only=True))
        sheets = wb.sheetnames
        target = sheet or sheets[0]
        ws = wb[target]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(c) if c is not None else "" for c in row])
        return {"path": path, "sheets": sheets, "active_sheet": target, "rows": len(rows), "columns": len(rows[0]) if rows else 0, "data": rows[:100]}
    except Exception as e:
        return {"error": str(e)}


async def create_excel(data: list[list[Any]], headers: list[str] | None = None, output_path: str | None = None) -> dict[str, Any]:
    if not HAS_OPENPYXL:
        return {"error": "openpyxl not installed"}
    try:
        wb = Workbook()
        ws = wb.active
        if headers:
            ws.append(headers)
        for row in data:
            ws.append(row)
        out = output_path or os.path.join(tempfile.gettempdir(), "friday_workbook.xlsx")
        await asyncio.get_event_loop().run_in_executor(None, lambda: wb.save(out))
        return {"path": out, "rows": len(data), "headers": headers}
    except Exception as e:
        return {"error": str(e)}


async def create_xlsx_chart(sections: list[dict[str, Any]], output_path: str | None = None) -> dict[str, Any]:
    """Create Excel with data sheets + chart image sheets."""
    if not HAS_OPENPYXL or not HAS_MPL:
        return {"error": "openpyxl or matplotlib not installed"}
    try:
        wb = Workbook()
        sheet_idx = 0
        for sec in sections:
            t = sec.get("type", "data")
            if t == "data":
                sheet_name = sec.get("name", f"Sheet{sheet_idx+1}")
                if sheet_idx == 0:
                    ws = wb.active
                    ws.title = sheet_name
                else:
                    ws = wb.create_sheet(title=sheet_name)
                sheet_idx += 1
                headers = sec.get("headers", [])
                if headers:
                    ws.append(headers)
                for row in sec.get("rows", []):
                    ws.append(row)
            elif t == "chart":
                chart_path = _make_chart(
                    sec.get("chart_type", "bar"), sec.get("data", []),
                    title=sec.get("title", ""),
                    xlabel=sec.get("xlabel", ""),
                    ylabel=sec.get("ylabel", ""),
                    labels=sec.get("labels"),
                    data2=sec.get("data2"),
                    data3=sec.get("data3"),
                )
                if chart_path:
                    from openpyxl.drawing.image import Image as XlImage
                    sheet_name = sec.get("name", f"Chart{sheet_idx+1}")
                    ws = wb.create_sheet(title=sheet_name)
                    sheet_idx += 1
                    img = XlImage(chart_path)
                    img.width = 600
                    img.height = 400
                    ws.add_image(img, "B2")
        out = output_path or os.path.join(tempfile.gettempdir(), "friday_workbook_charts.xlsx")
        await asyncio.get_event_loop().run_in_executor(None, lambda: wb.save(out))
        return {"path": out, "sections": len(sections)}
    except Exception as e:
        return {"error": str(e)}


# ── Data Frames (pandas) ──
HAS_PANDAS = False
try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    pass


async def analyze_csv(path: str) -> dict[str, Any]:
    if not HAS_PANDAS:
        return {"error": "pandas not installed"}
    if not os.path.exists(path):
        return {"error": f"File not found: {path}"}
    try:
        df = await asyncio.get_event_loop().run_in_executor(None, lambda: pd.read_csv(path))
        desc = df.describe(include="all").to_dict()
        return {"path": path, "rows": len(df), "columns": len(df.columns), "col_names": list(df.columns),
                "dtypes": {c: str(df[c].dtype) for c in df.columns}, "null_counts": df.isnull().sum().to_dict(),
                "summary": {c: {k: str(v) for k, v in desc[c].items()} for c in desc}}
    except Exception as e:
        return {"error": str(e)}


async def query_csv(path: str, query: str) -> dict[str, Any]:
    if not HAS_PANDAS:
        return {"error": "pandas not installed"}
    try:
        df = await asyncio.get_event_loop().run_in_executor(None, lambda: pd.read_csv(path))
        result = await asyncio.get_event_loop().run_in_executor(None, lambda: df.query(query))
        return {"query": query, "matches": len(result), "data": result.head(100).to_dict(orient="records") if len(result) > 0 else []}
    except Exception as e:
        return {"error": str(e)}


# ── PDF (pdfplumber / pypdf) ──
HAS_PDFPLUMBER = False
HAS_PYPDF = False
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    pass
try:
    import PyPDF2
    HAS_PYPDF = True
except ImportError:
    pass


async def read_pdf(path: str) -> dict[str, Any]:
    if not os.path.exists(path):
        return {"error": f"File not found: {path}"}
    if HAS_PDFPLUMBER:
        try:
            text = ""
            tables = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
                    tables.extend(page.extract_tables() or [])
            return {"path": path, "pages": len(pdf.pages), "text": text[:50000], "tables": len(tables)}
        except Exception as e:
            return {"error": str(e)}
    if HAS_PYPDF:
        try:
            text = ""
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
            return {"path": path, "pages": len(reader.pages), "text": text[:50000]}
        except Exception as e:
            return {"error": str(e)}
    return {"error": "pdfplumber or PyPDF2 not installed"}


# ── PDF Generation (reportlab + matplotlib) ──
HAS_REPORTLAB = False
HAS_MPL = False
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch, mm
    from reportlab.lib.colors import HexColor, black, white, navy, silver
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak,
        Table, TableStyle, Image, ListFlowable, ListItem,
        KeepTogether, FrameBreak, Flowable,
    )
    from reportlab.platypus.flowables import HRFlowable
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    HAS_REPORTLAB = True
except ImportError:
    pass
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mticker
    import numpy as np
    HAS_MPL = True
except ImportError:
    pass


def _make_chart(chart_type: str, data: list, title: str = "",
                xlabel: str = "", ylabel: str = "", labels: list | None = None,
                data2: list | None = None, data3: list | None = None) -> str | None:
    """Generate ANY chart type image. Returns file path or None."""
    if not HAS_MPL:
        return None
    try:
        colors = ["#2563eb", "#10b981", "#f59e0b", "#ef4444", "#8b5cf6",
                  "#ec4899", "#14b8a6", "#f97316", "#6366f1", "#84cc16",
                  "#06b6d4", "#d946ef", "#22c55e", "#eab308", "#a855f7"]

        is_3d = chart_type.startswith("3d_")
        if is_3d:
            from mpl_toolkits.mplot3d import Axes3D
            fig = plt.figure(figsize=(7, 5))
            ax = fig.add_subplot(111, projection="3d")
            fig.patch.set_facecolor("#f8f9fa")
            ax.set_facecolor("#f8f9fa")
        else:
            fig, ax = plt.subplots(figsize=(6, 3.5))
            fig.patch.set_facecolor("#f8f9fa")
            ax.set_facecolor("#f8f9fa")
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            ax.spines["left"].set_color("#cbd5e1")
            ax.spines["bottom"].set_color("#cbd5e1")
            ax.tick_params(colors="#475569", labelsize=9)

        # ── Bar ──
        if chart_type == "bar":
            x = range(len(data))
            ax.bar(x, data, color=colors[0], width=0.6, edgecolor="white", linewidth=0.5)
            if labels:
                ax.set_xticks(x)
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)

        # ── Horizontal Bar ──
        elif chart_type == "hbar":
            y = range(len(data))
            ax.barh(y, data, color=colors[0], height=0.6, edgecolor="white", linewidth=0.5)
            if labels:
                ax.set_yticks(y)
                ax.set_yticklabels(labels, fontsize=9)
            ax.invert_yaxis()

        # ── Grouped Bar ──
        elif chart_type == "grouped_bar":
            n = len(data)
            m = len(data2) if data2 else 0
            x = range(n)
            w = 0.35
            ax.bar([i - w/2 for i in x], data, width=w, color=colors[0], label="Series 1")
            if data2:
                ax.bar([i + w/2 for i in x], data2[:n], width=w, color=colors[1], label="Series 2")
            if labels:
                ax.set_xticks(x)
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
            ax.legend(fontsize=9)

        # ── Stacked Bar ──
        elif chart_type == "stacked_bar":
            x = range(len(data))
            ax.bar(x, data, color=colors[0], label="Series 1")
            if data2:
                ax.bar(x, data2[:len(data)], bottom=data, color=colors[1], label="Series 2")
            if labels:
                ax.set_xticks(x)
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
            ax.legend(fontsize=9)

        # ── Line ──
        elif chart_type == "line":
            ax.plot(data, color=colors[0], linewidth=2, marker="o", markersize=4)
            if data2:
                ax.plot(data2, color=colors[1], linewidth=2, marker="s", markersize=4, linestyle="--")
            if labels:
                ax.set_xticks(range(len(data)))
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)

        # ── Multi-line ──
        elif chart_type == "multi_line":
            series = [data]
            if data2: series.append(data2)
            if data3: series.append(data3)
            for i, s in enumerate(series):
                ax.plot(s, color=colors[i % len(colors)], linewidth=2, marker="o", markersize=3, label=f"Series {i+1}")
            if labels:
                ax.set_xticks(range(len(data)))
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
            ax.legend(fontsize=9)

        # ── Area ──
        elif chart_type == "area":
            x = range(len(data))
            ax.fill_between(x, data, alpha=0.3, color=colors[0])
            ax.plot(x, data, color=colors[0], linewidth=2)
            if data2:
                ax.fill_between(x, data2, alpha=0.3, color=colors[1])
                ax.plot(x, data2, color=colors[1], linewidth=2, linestyle="--")
            if labels:
                ax.set_xticks(x)
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)

        # ── Pie ──
        elif chart_type == "pie":
            wedges, texts, autotexts = ax.pie(
                data, labels=labels, autopct="%1.1f%%",
                colors=colors[:len(data)], startangle=90,
                textprops={"fontsize": 9}
            )
            for t in autotexts:
                t.set_color("white")
                t.set_fontweight("bold")

        # ── Donut ──
        elif chart_type == "donut":
            wedges, texts, autotexts = ax.pie(
                data, labels=labels, autopct="%1.1f%%",
                colors=colors[:len(data)], startangle=90,
                textprops={"fontsize": 9}, wedgeprops=dict(width=0.4)
            )
            for t in autotexts:
                t.set_color("white")
                t.set_fontweight("bold")

        # ── Scatter ──
        elif chart_type == "scatter":
            if data2:
                ax.scatter(data, data2, color=colors[0], s=50, alpha=0.7)
            else:
                ax.scatter(range(len(data)), data, color=colors[0], s=50, alpha=0.7)
            if labels and not data2:
                ax.set_xticks(range(len(data)))
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)

        # ── Bubble ──
        elif chart_type == "bubble":
            sizes = data3 or [30] * len(data)
            ax.scatter(data, data2 if data2 else range(len(data)), s=[max(10, s*2) for s in sizes],
                       color=colors[0], alpha=0.6, edgecolors="white", linewidth=0.5)

        # ── Histogram ──
        elif chart_type == "histogram":
            bins = data2 if data2 else 10
            if isinstance(bins, list):
                bins = len(bins) // 2 or 10
            ax.hist(data, bins=int(bins) if isinstance(bins, (int, float)) else 10,
                    color=colors[0], edgecolor="white", alpha=0.8)

        # ── Box ──
        elif chart_type == "box":
            box_data = [data]
            if data2: box_data.append(data2)
            if data3: box_data.append(data3)
            bp = ax.boxplot(box_data, patch_artist=True, widths=0.5)
            for patch, color in zip(bp["boxes"], colors[:len(box_data)]):
                patch.set_facecolor(color)
                patch.set_alpha(0.7)
            if labels:
                ax.set_xticklabels(labels[:len(box_data)], fontsize=9)

        # ── Violin ──
        elif chart_type == "violin":
            v_data = [data]
            if data2: v_data.append(data2)
            if data3: v_data.append(data3)
            vp = ax.violinplot(v_data, showmeans=True, showmedians=True)
            for i, body in enumerate(vp["bodies"]):
                body.set_facecolor(colors[i % len(colors)])
                body.set_alpha(0.7)
            if labels:
                ax.set_xticks(range(1, len(v_data) + 1))
                ax.set_xticklabels(labels[:len(v_data)], fontsize=9)

        # ── Heatmap ──
        elif chart_type == "heatmap":
            import numpy as np
            n = int(len(data) ** 0.5) or 1
            matrix = np.array(data[:n*n]).reshape(n, n)
            im = ax.imshow(matrix, cmap="Blues", aspect="auto")
            plt.colorbar(im, ax=ax, shrink=0.8)
            if labels:
                ax.set_xticks(range(n))
                ax.set_yticks(range(n))
                ax.set_xticklabels(labels[:n], rotation=30, ha="right", fontsize=8)
                ax.set_yticklabels(labels[:n], fontsize=8)

        # ── Radar ──
        elif chart_type == "radar":
            import numpy as np
            n = len(data)
            angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
            angles += angles[:1]
            values = data + data[:1]
            ax = fig.add_subplot(111, polar=True)
            ax.plot(angles, values, color=colors[0], linewidth=2)
            ax.fill(angles, values, alpha=0.25, color=colors[0])
            if labels:
                ax.set_xticks(angles[:-1])
                ax.set_xticklabels(labels, fontsize=9)

        # ── Candlestick ──
        elif chart_type == "candlestick":
            import numpy as np
            data_arr = np.array(data)
            if data_arr.ndim == 1:
                data_arr = data_arr.reshape(-1, 4) if len(data) % 4 == 0 else np.tile(data_arr, (4, 1)).T
            n = min(len(data_arr), 50)
            for i in range(n):
                row = data_arr[i]
                open_p, high, low, close = row[0], row[1], row[2], row[3]
                color = colors[1] if close >= open_p else colors[3]
                ax.plot([i, i], [low, high], color="#475569", linewidth=1)
                ax.plot([i - 0.2, i + 0.2], [open_p, open_p], color=color, linewidth=2.5)
                ax.plot([i - 0.2, i + 0.2], [close, close], color=color, linewidth=2.5)
            ax.set_xlim(-0.5, n - 0.5)

        # ── K-means clustering ──
        elif chart_type == "kmeans":
            from sklearn.cluster import KMeans
            import numpy as np
            n_clusters = data3[0] if data3 else 3
            x_vals = np.array(data)
            y_vals = np.array(data2) if data2 else np.random.RandomState(42).rand(len(data)) * 100
            X = np.column_stack([x_vals, y_vals])
            kmeans = KMeans(n_clusters=int(n_clusters), random_state=42, n_init=10)
            labels_pred = kmeans.fit_predict(X)
            centroids = kmeans.cluster_centers_
            cluster_colors = [colors[i % len(colors)] for i in range(int(n_clusters))]
            for i in range(int(n_clusters)):
                mask = labels_pred == i
                ax.scatter(X[mask, 0], X[mask, 1], c=cluster_colors[i], s=50,
                           alpha=0.7, label=f"Cluster {i+1}", edgecolors="white", linewidth=0.5)
            ax.scatter(centroids[:, 0], centroids[:, 1], c="black", s=200,
                       marker="X", label="Centroids", zorder=5)
            ax.legend(fontsize=8)

        # ── Contour ──
        elif chart_type == "contour":
            import numpy as np
            n = int(len(data) ** 0.5) or 10
            z = np.array(data[:n*n]).reshape(n, n)
            x = np.arange(n)
            y = np.arange(n)
            X, Y = np.meshgrid(x, y)
            cs = ax.contour(X, Y, z, levels=8, cmap="Blues")
            ax.clabel(cs, inline=True, fontsize=8)

        # ── 3D Scatter ──
        elif chart_type == "3d_scatter":
            import numpy as np
            z_vals = data3 if data3 else np.random.RandomState(42).rand(len(data)) * 100
            ax.scatter(data, data2 if data2 else range(len(data)), z_vals,
                       c=colors[0], s=30, alpha=0.7)
            ax.set_xlabel(xlabel or "X", fontsize=9, color="#475569")
            ax.set_ylabel(ylabel or "Y", fontsize=9, color="#475569")
            ax.set_zlabel("Z", fontsize=9, color="#475569")

        # ── 3D Surface ──
        elif chart_type == "3d_surface":
            import numpy as np
            n = int(len(data) ** 0.5) or 10
            z = np.array(data[:n*n]).reshape(n, n)
            x = np.arange(n)
            y = np.arange(n)
            X, Y = np.meshgrid(x, y)
            surf = ax.plot_surface(X, Y, z, cmap="Blues", edgecolor="none", alpha=0.8)
            fig.colorbar(surf, ax=ax, shrink=0.6)

        # ── 3D Bar ──
        elif chart_type == "3d_bar":
            import numpy as np
            x_pos = range(len(data))
            y_pos = data2[:len(data)] if data2 else [0] * len(data)
            z_pos = [0] * len(data)
            dx = dy = 0.5
            dz = data
            ax.bar3d(x_pos, y_pos, z_pos, dx, dy, dz, color=colors[0], alpha=0.8)
            if labels:
                ax.set_xticks(x_pos)
                ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=8)

        # ── Timeline ──
        elif chart_type == "timeline":
            import numpy as np
            events = labels if labels else [f"Event {i+1}" for i in range(len(data))]
            dates = data  # numeric positions (years, months, or sequential)
            colors_list = colors[:len(events)]
            y_pos = range(len(events))
            ax.set_yticks(list(y_pos))
            ax.set_yticklabels(events, fontsize=9)
            ax.set_xlabel(xlabel or "Timeline", fontsize=10, color="#475569")
            ax.set_title(title or "Timeline", fontsize=13, fontweight="bold", pad=12, color="#1e293b")
            ax.invert_yaxis()
            ax.spines["left"].set_visible(False)
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            ax.tick_params(left=False)
            for i, (d, c) in enumerate(zip(dates, colors_list)):
                ax.plot(d, i, marker="D", color=c, markersize=10, zorder=5)
                ax.axhline(y=i, xmin=0, xmax=1, color="#e2e8f0", linewidth=0.5, zorder=0)
            if data2:
                for i, (d, desc) in enumerate(zip(dates, data2)):
                    ax.annotate(str(desc), (d, i), textcoords="offset points",
                                xytext=(8, -4), fontsize=7, color="#475569",
                                ha="left", va="center")
            ax.set_xlim(min(dates) - (max(dates) - min(dates)) * 0.05 if len(dates) > 1 else dates[0] - 1,
                        max(dates) + (max(dates) - min(dates)) * 0.15 if len(dates) > 1 else dates[0] + 1)
            # Remove common styling override — already applied inline

        # Apply common styling
        if not is_3d:
            if chart_type != "timeline":
                if title:
                    ax.set_title(title, fontsize=13, fontweight="bold", pad=12, color="#1e293b")
                if xlabel:
                    ax.set_xlabel(xlabel, fontsize=10, color="#475569")
                if ylabel:
                    ax.set_ylabel(ylabel, fontsize=10, color="#475569")
                ax.yaxis.set_major_formatter(mticker.FormatStrFormatter("%d"))
            else:
                if ylabel:
                    ax.set_ylabel(ylabel, fontsize=10, color="#475569")
        else:
            if title:
                ax.set_title(title, fontsize=13, fontweight="bold", pad=12, color="#1e293b")

        plt.tight_layout()
        out_path = os.path.join(tempfile.gettempdir(),
                                f"friday_chart_{abs(hash(str(data)+str(data2)+title+chart_type))}.png")
        fig.savefig(out_path, dpi=150, bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig)
        return out_path
    except Exception:
        return None


async def create_pdf(
    sections: list[dict[str, Any]],
    title: str = "Friday Report",
    output_path: str | None = None,
) -> dict[str, Any]:
    """
    Generate a rich PDF with headings, paragraphs, tables, charts, and lists.
    IMPORTANT: First read friday/skills/pdf/SKILL.md for the usage guide and patterns.

    Each section dict supports:
      {"type": "heading", "text": "...", "level": 1|2|3}
      {"type": "paragraph", "text": "..."}
      {"type": "table", "headers": [...], "rows": [[...], ...], "caption": "..."}
      {"type": "chart", "chart_type": "bar"|"line"|"pie"|"hbar"|"scatter"|"timeline"|"area"|"multi_line"|"grouped_bar"|"stacked_bar"|"donut"|"bubble"|"histogram"|"box"|"violin"|"heatmap"|"radar"|"candlestick"|"kmeans"|"contour"|"3d_scatter"|"3d_surface"|"3d_bar",
                    "data": [...], "labels": [...], "title": "...",
                    "xlabel": "...", "ylabel": "..."}
      {"type": "bullets", "items": ["...", ...]}
      {"type": "numbered", "items": ["...", ...]}
      {"type": "divider"}
    """
    if not HAS_REPORTLAB:
        return {"error": "reportlab not installed. Install: pip install reportlab"}

    try:
        out = output_path or os.path.join(tempfile.gettempdir(), "friday_report.pdf")
        doc = SimpleDocTemplate(
            out, pagesize=letter,
            topMargin=0.75 * inch, bottomMargin=0.75 * inch,
            leftMargin=0.75 * inch, rightMargin=0.75 * inch,
        )

        styles = getSampleStyleSheet()
        s_normal = ParagraphStyle("CustomNormal", parent=styles["Normal"],
                                  fontSize=10, leading=14, spaceAfter=6,
                                  textColor=HexColor("#1e293b"))
        s_h1 = ParagraphStyle("H1", parent=styles["Heading1"],
                              fontSize=22, leading=28, spaceAfter=14, spaceBefore=8,
                              textColor=HexColor("#1e293b"), fontName="Helvetica-Bold")
        s_h2 = ParagraphStyle("H2", parent=styles["Heading2"],
                              fontSize=16, leading=20, spaceAfter=10, spaceBefore=12,
                              textColor=HexColor("#334155"), fontName="Helvetica-Bold")
        s_h3 = ParagraphStyle("H3", parent=styles["Heading3"],
                              fontSize=13, leading=17, spaceAfter=8, spaceBefore=8,
                              textColor=HexColor("#475569"), fontName="Helvetica-Bold")
        s_title = ParagraphStyle("Title", parent=styles["Title"],
                                 fontSize=28, leading=34, spaceAfter=6,
                                 textColor=HexColor("#0f172a"), alignment=TA_CENTER,
                                 fontName="Helvetica-Bold")
        s_subtitle = ParagraphStyle("Subtitle", parent=styles["Normal"],
                                    fontSize=12, leading=16, spaceAfter=24,
                                    textColor=HexColor("#64748b"), alignment=TA_CENTER)
        s_caption = ParagraphStyle("Caption", parent=styles["Normal"],
                                   fontSize=9, leading=12, spaceAfter=10,
                                   textColor=HexColor("#64748b"), alignment=TA_LEFT)
        s_bullet = ParagraphStyle("Bullet", parent=s_normal,
                                  leftIndent=20, bulletIndent=8, spaceBefore=2, spaceAfter=2)
        s_code = ParagraphStyle("Code", parent=styles["Code"],
                                fontName="Courier", fontSize=8, leading=11,
                                backColor=HexColor("#f1f5f9"),
                                leftIndent=10, rightIndent=10, spaceAfter=8,
                                borderPadding=6)

        story = []

        # ---- Title Page ----
        story.append(Spacer(1, 1.5 * inch))
        story.append(Paragraph(title, s_title))
        story.append(Spacer(1, 0.3 * inch))
        story.append(HRFlowable(width="60%", thickness=2, color=HexColor("#2563eb"),
                                spaceAfter=12, spaceBefore=6, hAlign="CENTER"))
        story.append(Paragraph("Generated by FRIDAY · Stark Industries AI", s_subtitle))
        story.append(PageBreak())

        def _add_paragraph(text: str, style: ParagraphStyle = s_normal) -> None:
            for para in text.split("\n\n"):
                para = para.strip()
                if para:
                    story.append(Paragraph(para.replace("\n", "<br/>"), style))

        # ---- Content Sections ----
        for sec in sections:
            sec_type = sec.get("type", "paragraph")

            if sec_type == "heading":
                level = sec.get("level", 1)
                style_map = {1: s_h1, 2: s_h2, 3: s_h3}
                hs = style_map.get(level, s_h2)
                story.append(Paragraph(sec["text"], hs))

            elif sec_type == "paragraph":
                _add_paragraph(sec["text"], s_normal)

            elif sec_type == "bullets":
                items = sec.get("items", [])
                for item in items:
                    story.append(Paragraph(f"• {item}", s_bullet))
                story.append(Spacer(1, 4))

            elif sec_type == "numbered":
                items = sec.get("items", [])
                for idx, item in enumerate(items, 1):
                    story.append(Paragraph(f"{idx}. {item}", s_bullet))
                story.append(Spacer(1, 4))

            elif sec_type == "table":
                headers = sec.get("headers", [])
                rows = sec.get("rows", [])
                caption = sec.get("caption", "")
                if rows:
                    table_data = [headers] + rows if headers else rows
                    col_w = (doc.width) / max(len(table_data[0]), 1)
                    # Style header row separately
                    tbl = Table(table_data, colWidths=[col_w] * len(table_data[0]),
                                repeatRows=1 if headers else 0)
                    tbl_style_cmds = [
                        ("FONTSIZE", (0, 0), (-1, -1), 9),
                        ("LEADING", (0, 0), (-1, -1), 12),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#cbd5e1")),
                    ]
                    if headers:
                        tbl_style_cmds += [
                            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1e293b")),
                            ("TEXTCOLOR", (0, 0), (-1, 0), white),
                            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                            ("FONTSIZE", (0, 0), (-1, 0), 10),
                        ]
                        # Alternate row colors
                        for ri in range(1, len(table_data)):
                            if ri % 2 == 0:
                                tbl_style_cmds.append(
                                    ("BACKGROUND", (0, ri), (-1, ri), HexColor("#f8fafc"))
                                )
                    tbl.setStyle(TableStyle(tbl_style_cmds))
                    story.append(tbl)
                    if caption:
                        story.append(Paragraph(f"<i>{caption}</i>", s_caption))
                    story.append(Spacer(1, 8))

            elif sec_type == "chart":
                chart_type = sec.get("chart_type", "bar")
                data = sec.get("data", [])
                if data and HAS_MPL:
                    chart_path = _make_chart(
                        chart_type, data,
                        title=sec.get("title", ""),
                        xlabel=sec.get("xlabel", ""),
                        ylabel=sec.get("ylabel", ""),
                        labels=sec.get("labels"),
                        data2=sec.get("data2"),
                        data3=sec.get("data3"),
                    )
                    if chart_path:
                        img = Image(chart_path, width=6.5 * inch, height=3.8 * inch)
                        story.append(img)
                        story.append(Spacer(1, 8))
                elif data and not HAS_MPL:
                    # Fallback: render data as text
                    story.append(Paragraph(f"<b>Chart ({chart_type}):</b> {sec.get('title', '')}", s_normal))
                    story.append(Paragraph(f"Data: {', '.join(str(d) for d in data)}", s_caption))

            elif sec_type == "divider":
                story.append(Spacer(1, 6))
                story.append(HRFlowable(width="100%", thickness=0.5, color=HexColor("#e2e8f0"),
                                        spaceAfter=6, spaceBefore=6))
                story.append(Spacer(1, 6))

            elif sec_type == "image":
                img_path = sec.get("path", "")
                if img_path and os.path.exists(img_path):
                    img = Image(img_path, width=min(6.5 * inch, doc.width),
                                height=min(4 * inch, doc.height * 0.6))
                    story.append(img)
                    story.append(Spacer(1, 8))

            elif sec_type == "code":
                code_text = sec.get("text", "")
                story.append(Paragraph(code_text.replace("\n", "<br/>"), s_code))
                story.append(Spacer(1, 4))

        await asyncio.get_event_loop().run_in_executor(None, lambda: doc.build(story))
        section_count = len(sections)
        return {"path": out, "sections": section_count, "title": title}
    except Exception as e:
        return {"error": str(e)}


# ── PowerPoint (python-pptx) ──
HAS_PPTX = False
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    pass


async def read_pptx(path: str) -> dict[str, Any]:
    if not HAS_PPTX:
        return {"error": "python-pptx not installed"}
    if not os.path.exists(path):
        return {"error": f"File not found: {path}"}
    try:
        prs = await asyncio.get_event_loop().run_in_executor(None, lambda: Presentation(path))
        slides = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        slide_text.append(para.text)
            slides.append({"text": "\n".join(slide_text)[:500]})
        return {"path": path, "slides": len(slides), "slide_width": prs.slide_width, "slide_height": prs.slide_height,
                "content": slides[:50]}
    except Exception as e:
        return {"error": str(e)}


async def create_pptx(title: str, slides: list[dict[str, Any]], output_path: str | None = None) -> dict[str, Any]:
    """Create PowerPoint with text and chart slides."""
    if not HAS_PPTX:
        return {"error": "python-pptx not installed"}
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_slide.shapes.title.text = title
        for s in slides:
            slide_type = s.get("type", "content")
            if slide_type == "chart":
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.title.text = s.get("title", "")
                chart_path = _make_chart(
                    s.get("chart_type", "bar"), s.get("data", []),
                    title=s.get("title", ""),
                    xlabel=s.get("xlabel", ""),
                    ylabel=s.get("ylabel", ""),
                    labels=s.get("labels"),
                    data2=s.get("data2"),
                    data3=s.get("data3"),
                )
                if chart_path:
                    from pptx.util import Inches as PptInches
                    slide.shapes.add_picture(chart_path, PptInches(0.5), PptInches(1.2),
                                             width=PptInches(12), height=PptInches(5.5))
                if s.get("notes"):
                    slide.notes_slide.notes_text_frame.text = s["notes"]
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = s.get("title", "")
                slide.placeholders[1].text = s.get("content", "")
                if s.get("notes"):
                    slide.notes_slide.notes_text_frame.text = s["notes"]
        out = output_path or os.path.join(tempfile.gettempdir(), "friday_presentation.pptx")
        await asyncio.get_event_loop().run_in_executor(None, lambda: prs.save(out))
        return {"path": out, "slides": len(slides) + 1}
    except Exception as e:
        return {"error": str(e)}
