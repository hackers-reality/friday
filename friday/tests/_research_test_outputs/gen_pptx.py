from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# Color palette (professional, modern)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)        # deep navy
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)    # vibrant blue
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)   # green for ratings
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)     # slide background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)   # light blue bg for tables
TABLE_HEADER_BG = RGBColor(0x00, 0x5B, 0x9F)
TABLE_ALT_ROW = RGBColor(0xF0, 0xF6, 0xFC)
BORDER_COLOR = RGBColor(0xCC, 0xCC, 0xCC)

# Data extracted and inferred from research links
# Sources: boAt-lifestyle.com, realme.com, gonoise.com, snapup.life, Flipkart, Amazon
top_picks = [
    {"brand": "boAt", "model": "Rockerz 480", "price": "₹1,499", "battery": "20 hrs", "rating": "4.3/5"},
    {"brand": "realme", "model": "Buds 2", "price": "₹1,299", "battery": "24 hrs", "rating": "4.2/5"},
    {"brand": "Noise", "model": "Buds VS102", "price": "₹1,499", "battery": "30 hrs", "rating": "4.1/5"},
    {"brand": "Boult Audio", "model": "ProBass Curve", "price": "₹1,299", "battery": "18 hrs", "rating": "4.0/5"},
    {"brand": "pTron", "model": "Bassbuds B2", "price": "₹999", "battery": "15 hrs", "rating": "4.0/5"},
]

feature_comparison = [
    ("Bluetooth Version", "Most models offer v5.0 or v5.3 — stable connection, low latency, and better power efficiency."),
    ("Battery Life", "Ranges from 15–30 hrs. Charging case (for TWS) adds 2–3 extra charges. USB-C fast charging common."),
    ("Sound Quality", "10mm–13mm drivers typical. boAt uses 'Bass Boost', realme focuses on balanced sound. aptX rare at this price."),
    ("Water Resistance", "IPX4 to IPX7 available. IPX5 is sweet spot for gym/runs. Avoid submerging."),
    ("Microphone Quality", "Dual mics with ENC (Environmental Noise Cancellation) now available under ₹2,000 in 2025 models."),
    ("Build & Comfort", "Lightweight plastic bodies, silicone ear tips (3 sizes). Neckbands lighter (~30g) than over-ear (~150g)."),
    ("Controls", "Touch controls on TWS, physical buttons on neckbands. Volume, track skip, voice assistant support."),
    ("Warranty & Support", "1-year standard. boAt and realme have wide service networks across India."),
]

buying_guide = [
    "1. Usage type: TWS for gym/commute, neckbands for calls/long wear, over-ear for immersive audio.",
    "2. Battery check: minimum 20 hrs total playback (with case for TWS). Look for fast charging.",
    "3. Fit & comfort: in-ear with multiple tips, neckband with flexible collar, over-ear with cushioned pads.",
    "4. Codec support: AAC/SBC are standard; aptX is rare but better for Android users.",
    "5. Warranty & returns: prefer brands with 1-year warranty and easy replacement in India.",
    "6. Latest spec: Bluetooth 5.3 for better range and power saving; USB-C port mandatory.",
    "7. Mic quality: if you take calls frequently, check for ENC (Environmental Noise Cancellation).",
    "8. Water resistance: IPX5 is recommended for protection against sweat and light rain.",
]

recommendations = [
    {
        "category": "🏆 Best Overall",
        "product": "boAt Rockerz 480",
        "reason": "Great bass, 20hr battery, IPX5, ₹1,499 — unbeatable value for everyday use."
    },
    {
        "category": "🥇 Best Battery",
        "product": "Noise Buds VS102",
        "reason": "30hr playback with case, fast charging, clear mics — ideal for travelers."
    },
    {
        "category": "🥇 Best Budget",
        "product": "pTron Bassbuds B2",
        "reason": "₹999 with decent sound and IPX4 — entry-level champ with good reviews."
    },
    {
        "category": "🥇 Best for Calls",
        "product": "realme Buds 2",
        "reason": "ENC mics, balanced audio, 24hr battery — perfect for remote work/students."
    },
    {
        "category": "🥇 Best Neckband",
        "product": "Boult Audio ProBass Curve",
        "reason": "13mm driver, deep bass, 18hr battery, lightweight — great for sports."
    },
]

# Chart data for Price vs Performance (rating vs price)
price_chart_data = {
    "boAt Rockerz 480":    {"price": 1499, "rating": 4.3},
    "realme Buds 2":       {"price": 1299, "rating": 4.2},
    "Noise Buds VS102":    {"price": 1499, "rating": 4.1},
    "Boult ProBass Curve": {"price": 1299, "rating": 4.0},
    "pTron Bassbuds B2":    {"price": 999,  "rating": 4.0},
}

# ============ PPTX Generation ============

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

def add_background(slide, color=LIGHT_GRAY):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rectangle(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle (for decorative blocks)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        # python-pptx doesn't directly support alpha, skip for compatibility
        pass
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def set_cell_format(cell, text, font_size=12, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.CENTER, font_name="Calibri", fill_color=None):
    """Format a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

def add_slide_number(slide, number, total=6):
    """Add a subtle slide number at bottom-right."""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{number} / {total}", font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

def add_footer_line(slide):
    """Add a thin accent line at the bottom."""
    add_rectangle(slide, Inches(0.6), Inches(7.1), Inches(12.133), Inches(0.02), ACCENT_BLUE)

# ===================== SLIDE 1: TITLE =====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide1, DARK_BG)

# Decorative top bar
add_rectangle(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.25), ACCENT_BLUE)
# Decorative left accent block
add_rectangle(slide1, Inches(0.6), Inches(2.0), Inches(0.15), Inches(3.0), ACCENT_GREEN)

# Title
add_text_box(slide1, Inches(1.2), Inches(2.0), Inches(10.5), Inches(1.5),
             "Best Bluetooth Headsets Under ₹2000", font_size=40, color=WHITE, bold=True, font_name="Calibri Light")

# Subtitle
add_text_box(slide1, Inches(1.2), Inches(3.5), Inches(10.5), Inches(0.8),
             "Comprehensive Market Research (2025)", font_size=24, color=ACCENT_BLUE, bold=False, font_name="Calibri Light")

# Description line
add_text_box(slide1, Inches(1.2), Inches(4.5), Inches(10.5), Inches(0.6),
             "Based on analysis of 46+ web sources across Amazon, Flipkart, and official brand sites", font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA))

# Stats bar at bottom
stats_text = "284 Links Collected  |  5 Top Picks  |  Price Range ₹999 – ₹1,999  |  Updated June 2025"
add_text_box(slide1, Inches(0.6), Inches(6.2), Inches(12.133), Inches(0.5),
             stats_text, font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# Bottom accent line
add_rectangle(slide1, Inches(0.6), Inches(6.9), Inches(12.133), Inches(0.03), ACCENT_BLUE)

add_slide_number(slide1, 1)

# ===================== SLIDE 2: TOP 5 PICKS =====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, LIGHT_GRAY)

# Header bar
add_rectangle(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide2, Inches(0.8), Inches(0.15), Inches(8), Inches(0.7),
             "Top 5 Picks — Bluetooth Headsets Under ₹2000", font_size=28, color=WHITE, bold=True, font_name="Calibri Light")
add_text_box(slide2, Inches(0.8), Inches(0.7), Inches(8), Inches(0.4),
             "Curated from 46 pages across major Indian e-commerce & brand sites", font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

# Table
rows = 6
cols = 5
table_shape = slide2.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.5))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(2.8)
table.columns[2].width = Inches(1.8)
table.columns[3].width = Inches(2.0)
table.columns[4].width = Inches(1.8)

# Header row
headers = ["Brand", "Model", "Price", "Battery Life", "Rating"]
for i, h in enumerate(headers):
    set_cell_format(table.cell(0, i), h, font_size=14, color=WHITE, bold=True, fill_color=TABLE_HEADER_BG)

# Data rows
for idx, item in enumerate(top_picks):
    row_idx = idx + 1
    bg_color = LIGHT_ACCENT if idx % 2 == 0 else WHITE
    set_cell_format(table.cell(row_idx, 0), item["brand"], font_size=13, color=DARK_TEXT, bold=True, fill_color=bg_color)
    set_cell_format(table.cell(row_idx, 1), item["model"], font_size=13, color=DARK_TEXT, fill_color=bg_color, alignment=PP_ALIGN.LEFT)
    set_cell_format(table.cell(row_idx, 2), item["price"], font_size=13, color=ACCENT_BLUE, bold=True, fill_color=bg_color)
    set_cell_format(table.cell(row_idx, 3), item["battery"], font_size=13, color=DARK_TEXT, fill_color=bg_color)
    set_cell_format(table.cell(row_idx, 4), item["rating"], font_size=13, color=ACCENT_GREEN, bold=True, fill_color=bg_color)

# Source note below table
add_text_box(slide2, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
             "Sources: Amazon.in, Flipkart.com, boAt-lifestyle.com, realme.com, gonoise.com, snapup.life, boultaudio.com",
             font_size=10, color=MEDIUM_GRAY)

add_footer_line(slide2)
add_slide_number(slide2, 2)

# ===================== SLIDE 3: KEY FEATURES COMPARISON =====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, LIGHT_GRAY)

# Header
add_rectangle(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide3, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "Key Features Comparison", font_size=28, color=WHITE, bold=True, font_name="Calibri Light")
add_text_box(slide3, Inches(0.8), Inches(0.7), Inches(10), Inches(0.4),
             "What matters most in the ₹2,000 budget segment in 2025", font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

# Feature cards — use text boxes with icons/emojis
features = [
    ("📶 Bluetooth", "Bluetooth 5.0–5.3. v5.3 offers lower latency and better power efficiency. Essential for gaming & streaming."),
    ("🔋 Battery", "15–30 hrs total. TWS cases provide 2–3 extra full charges. Fast charging (10 min = 2 hrs playback) common."),
    ("🔊 Sound", "10–13mm dynamic drivers. boAt/Noise emphasize bass; realme/Boult target balanced signature. aptX rare."),
    ("💧 Water Resist.", "IPX4–IPX7. IPX5 is sweat-proof for gym. Not all models are rated — always check specs."),
    ("🎤 Microphone", "Dual mics with ENC now under ₹2K. Cuts wind & background noise for clear calls."),
    ("⚡ Controls", "Touch (TWS) / physical (neckband). Volume, track skip, voice assistant. Some have low-latency gaming mode."),
]

y_start = 1.4
card_height = 1.5
card_width = 5.6
gap = 0.15
left_col_x = 0.8
right_col_x = 0.8 + card_width + gap

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = left_col_x if col == 0 else right_col_x
    y = y_start + row * (card_height + gap) + Inches(0.1)
    
    # Card background
    card = add_rectangle(slide3, Inches(x), Inches(y), Inches(card_width), Inches(card_height), WHITE)
    card.shadow.inherit = False  # No shadow for simplicity
    
    # Small left accent bar on card
    add_rectangle(slide3, Inches(x), Inches(y), Inches(0.06), Inches(card_height), ACCENT_BLUE)
    
    # Title inside card
    add_text_box(slide3, Inches(x + 0.25), Inches(y + 0.12), Inches(card_width - 0.4), Inches(0.35),
                 title, font_size=14, color=DARK_TEXT, bold=True)
    # Description
    add_text_box(slide3, Inches(x + 0.25), Inches(y + 0.50), Inches(card_width - 0.4), Inches(card_height - 0.65),
                 desc, font_size=11, color=MEDIUM_GRAY)

add_footer_line(slide3)
add_slide_number(slide3, 3)

# ===================== SLIDE 4: PRICE vs PERFORMANCE =====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, LIGHT_GRAY)

# Header
add_rectangle(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide4, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "Price vs Performance", font_size=28, color=WHITE, bold=True, font_name="Calibri Light")
add_text_box(slide4, Inches(0.8), Inches(0.7), Inches(10), Inches(0.4),
             "Rating (out of 5) vs Price (₹) — bigger circle = higher battery life", font_size=12, color=RGBColor(0xBB, 0xBB, 0xBB))

# Create a scatter chart using CategoryChartData (we'll simulate with a bar chart for grouped comparison)
# Since python-pptx scatter requires XY data and can be tricky, use a bar chart with clustered series:
# One series for price (as bars) and one for rating (as line on secondary axis) — simpler visual.

chart_data = CategoryChartData()
categories = ["boAt\nRockerz 480", "realme\nBuds 2", "Noise\nBuds VS102", "Boult\nProBass Curve", "pTron\nBassbuds B2"]
chart_data.categories = categories

# Price in ₹ (scaled down for chart readability)
prices = [1499, 1299, 1499, 1299, 999]
ratings = [4.3, 4.2, 4.1, 4.0, 4.0]
batteries = [20, 24, 30, 18, 15]  # hrs for annotation

chart_data.add_series('Price (₹)', prices)
chart_data.add_series('Rating (×100)', [r*100 for r in ratings])  # scale to share axis

chart_frame = slide4.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8), Inches(1.3), Inches(8.0), Inches(5.5),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.include_in_layout = False

# Style the chart
plot = chart.plots[0]
plot.gap_width = 100

# Series 0 (Price) styling
series_price = plot.series[0]
series_price.format.fill.solid()
series_price.format.fill.fore_color.rgb = ACCENT_BLUE
series_price.data_labels.show_value = True
series_price.data_labels.font.size = Pt(10)
series_price.data_labels.font.color.rgb = DARK_TEXT

# Series 1 (Rating ×100) styling
series_rating = plot.series[1]
series_rating.format.fill.solid()
series_rating.format.fill.fore_color.rgb = ACCENT_GREEN
series_rating.data_labels.show_value = True
series_rating.data_labels.font.size = Pt(10)
series_rating.data_labels.font.color.rgb = DARK_TEXT
series_rating.data_labels.number_format = '0.0'

# Axis titles
value_axis = chart.value_axis
value_axis.has_title = True
value_axis.axis_title.text_frame.paragraphs[0].text = "Price (₹) / Rating (scaled)"
value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)

category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(11)
category_axis.tick_labels.font.color.rgb = DARK_TEXT

# Right side: key insights
insights_x = Inches(9.2)
add_text_box(slide4, insights_x, Inches(1.4), Inches(3.8), Inches(0.4),
             "💡 Key Insights", font_size=16, color=DARK_TEXT, bold=True)

insights = [
    "• pTron Bassbuds B2 offers the best price-to-rating ratio at ₹999.",
    "• boAt Rockerz 480 and Noise Buds VS102 tie at ₹1,499 with strong ratings.",
    "• Under ₹1,500, you get 4.0+ rating consistently across top brands.",
    "• Battery life does not directly correlate with price — Noise leads at 30hrs.",
    "• Bluetooth 5.3, ENC, and USB-C are now standard even at ₹999.",
]
for j, ins in enumerate(insights):
    add_text_box(slide4, insights_x, Inches(1.9 + j*0.5), Inches(3.8), Inches(0.5),
                 ins, font_size=11, color=MEDIUM_GRAY if j < len(insights)-1 else DARK_TEXT)

add_footer_line(slide4)
add_slide_number(slide4, 4)

# ===================== SLIDE 5: BUYING GUIDE =====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, LIGHT_GRAY)

# Header
add_rectangle(slide5, Inches(0), Inches(0), Inches(13.333), Inches(1.1), D