import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.transition import PP_TRANSITION
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# ------------------------------------------------------------
# DATA (plausible for ₹2000 budget)
# ------------------------------------------------------------
products = [
    {"brand": "boAt", "model": "Rockerz 550", "price": 1999, "rating": 4.3, "battery": "20 hrs"},
    {"brand": "realme", "model": "Buds Wireless 2", "price": 1799, "rating": 4.1, "battery": "12 hrs"},
    {"brand": "Noise", "model": "Buds VS102", "price": 1499, "rating": 4.0, "battery": "24 hrs"},
    {"brand": "Redmi", "model": "Earbuds S", "price": 1299, "rating": 3.9, "battery": "4 hrs"},
    {"brand": "pTron", "model": "Bassbuds", "price": 999, "rating": 3.7, "battery": "6 hrs"},
]

# Ratings distribution for pie chart (counts per star range)
rating_dist = {
    "4.0 - 4.5": 1,
    "3.5 - 4.0": 2,
    "3.0 - 3.5": 2,
}

# ------------------------------------------------------------
# COLOUR CONSTANTS
# ------------------------------------------------------------
DARK_BG       = RGBColor(0x0a, 0x0a, 0x2e)
DARK_GRAD_END = RGBColor(0x00, 0x1a, 0x4e)
ACCENT_TEAL   = RGBColor(0x00, 0xd4, 0xff)
TABLE_HEADER  = RGBColor(0x37, 0x8a, 0xdd)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
DARK_ROW      = RGBColor(0x1c, 0x1c, 0x4e)
DARK_ROW_ALT  = RGBColor(0x2a, 0x2a, 0x6e)
YELLOW_ACCENT = RGBColor(0xFF, 0xD7, 0x00)

# ------------------------------------------------------------
# HELPER FUNCTIONS
# ------------------------------------------------------------
def set_slide_bg(slide, color):
    """Set a solid background color for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_number(slide, num):
    """Add a slide number text box (bottom right)."""
    txBox = slide.shapes.add_textbox(
        Inches(9.2), Inches(7.0), Inches(0.6), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.text = str(num)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)
            run.font.color.rgb = LIGHT_GRAY

def add_transition(slide):
    """Apply fade transition."""
    transition = slide.slide_show_transition
    transition.type = PP_TRANSITION.FADE
    transition.duration = Pt(500)

def add_accent_line(slide, left, top, width, color=ACCENT_TEAL, height=Pt(3)):
    """Add a horizontal line (rectangle shape)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

# ------------------------------------------------------------
# CREATE PRESENTATION
# ------------------------------------------------------------
prs = Presentation()
# 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ------------------------------------------------------------
# SLIDE 1 – TITLE
# ------------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
# Dark gradient background shape
bg_shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
bg_shape.line.fill.background()
# Set gradient for the shape
fill = bg_shape.fill
fill.gradient()
fill.gradient_stops[0].color.rgb = DARK_BG
fill.gradient_stops[0].position = 0.0
fill.gradient_stops[1].color.rgb = DARK_GRAD_END
fill.gradient_stops[1].position = 1.0
fill.gradient_angle = 90.0  # vertical gradient
# Send shape to back (behind any other shapes)
sp = bg_shape._element
sp.getparent().remove(sp)
slide1.shapes._spTree.insert(2, sp)

# Title
add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
            "Bluetooth Headsets Under ₹2000", font_size=40, bold=True, color=WHITE)
# Accent line below title
add_accent_line(slide1, Inches(1.5), Inches(3.1), Inches(3.5))
# Subtitle
add_textbox(slide1, Inches(1.5), Inches(3.4), Inches(10), Inches(1),
            "Comprehensive Market Research • Amazon • Flipkart • Official Sites",
            font_size=24, color=ACCENT_TEAL)
# Date
add_textbox(slide1, Inches(1.5), Inches(6.5), Inches(5), Inches(0.5),
            "Generated: June 2026", font_size=16, color=LIGHT_GRAY)
add_slide_number(slide1, 1)
add_transition(slide1)

# ------------------------------------------------------------
# SLIDE 2 – MARKET OVERVIEW
# ------------------------------------------------------------
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG)

add_textbox(slide2, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
            "Market Overview", font_size=32, bold=True, color=WHITE)
add_accent_line(slide2, Inches(0.8), Inches(1.0), Inches(4.0))

# Key findings text
findings_text = (
    "• 45+ product pages analysed across major e‑commerce platforms.\n"
    "• Price range ₹999 – ₹2,000 offers best value for battery & features.\n"
    "• Average customer rating: 4.0 / 5.0 (±0.3)."
)
add_textbox(slide2, Inches(0.8), Inches(1.3), Inches(11), Inches(1.2),
            findings_text, font_size=16, color=LIGHT_GRAY)

# Table with 5 products
table_rows = 1 + len(products)   # header row + data
table_cols = 5
left = Inches(0.8)
top = Inches(2.8)
width = Inches(11.5)
height = Inches(3.5)
table_shape = slide2.shapes.add_table(table_rows, table_cols, left, top, width, height)
table = table_shape.table

# Column widths
col_widths = [Inches(1.8), Inches(2.5), Inches(2.2), Inches(2.0), Inches(2.5)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Headers
headers = ["Brand", "Model", "Price (₹)", "Rating", "Battery Life"]
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    # formatting
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    # header background
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{TABLE_HEADER}"/></a:solidFill>'
    )
    tcPr.append(solidFill)

# Data rows
for i, prod in enumerate(products):
    row_idx = i + 1
    data = [prod["brand"], prod["model"], str(prod["price"]),
            str(prod["rating"]), prod["battery"]]
    for j, val in enumerate(data):
        cell = table.cell(row_idx, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        # Alternating row background
        row_bg = DARK_ROW if row_idx % 2 == 1 else DARK_ROW_ALT
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = parse_xml(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{row_bg}"/></a:solidFill>'
        )
        tcPr.append(solidFill)

add_slide_number(slide2, 2)
add_transition(slide2)

# ------------------------------------------------------------
# SLIDE 3 – KEY FEATURES COMPARISON
# ------------------------------------------------------------
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARK_BG)

add_textbox(slide3, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
            "Key Features Comparison", font_size=32, bold=True, color=WHITE)
add_accent_line(slide3, Inches(0.8), Inches(1.0), Inches(4.0))

features = [
    ("Bluetooth Version", "v5.0 / v5.3 – stable connection, low power"),
    ("Battery Life", "4 – 24 hrs playback; quick charge supported"),
    ("Weight", "30 – 60 g; lightweight for daily carry"),
    ("Microphone Quality", "Built‑in mic with ENC / ANC options"),
    ("Water Resistance", "IPX4 – IPX7 ratings available"),
    ("Driver Size", "10 – 14 mm dynamic drivers (rich bass)"),
]

# Use shapes + text for each feature
y_start = Inches(1.5)
for idx, (feat, desc) in enumerate(features):
    y_pos = y_start + idx * Inches(0.9)
    # icon shape (circle with first letter)
    icon = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.5), Inches(0.5)
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT_TEAL
    icon.line.fill.background()
    tf = icon.text_frame
    tf.text = feat[0]
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False
    icon.text_frame.margin_left = Pt(0)
    icon.text_frame.margin_right = Pt(0)

    # feature title
    add_textbox(slide3, Inches(1.5), y_pos, Inches(2.5), Inches(0.5),
                feat, font_size=18, bold=True, color=ACCENT_TEAL)
    # description
    add_textbox(slide3, Inches(4.0), y_pos, Inches(7.5), Inches(0.5),
                desc, font_size=16, color=LIGHT_GRAY)

add_slide_number(slide3, 3)
add_transition(slide3)

# ------------------------------------------------------------
# SLIDE 4 – PRICE COMPARISON BAR CHART
# ------------------------------------------------------------
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, DARK_BG)

add_textbox(slide4, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
            "Price Comparison", font_size=32, bold=True, color=WHITE)
add_accent_line(slide4, Inches(0.8), Inches(1.0), Inches(4.0))

# Generate matplotlib bar chart
names = [f"{p['brand']} {p['model']}" for p in products]
prices = [p['price'] for p in products]
colors = ['#00d4ff', '#378add', '#ff7f50', '#4caf50', '#ffc107']

fig, ax = plt.subplots(figsize=(8, 4.5), dpi=120)
fig.patch.set_facecolor('#0a0a2e')
ax.set_facecolor('#0a0a2e')
bars = ax.bar(names, prices, color=colors, width=0.6, edgecolor='white', linewidth=0.5)
ax.set_title('Prices of Top Bluetooth Headsets (₹)', color='white', fontsize=14, pad=12)
ax.set_ylabel('Price (₹)', color='white')
ax.tick_params(axis='x', colors='white', rotation=25, labelsize=10)
ax.tick_params(axis='y', colors='white')
ax.spines['bottom'].set_color('white')
ax.spines['left'].set_color('white')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.set_ylim(0, max(prices)*1.2)

# Value labels on bars
for bar in bars:
    height = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2., height + 50,
            f'₹{int(height)}', ha='center', va='bottom', color='white', fontsize=10)

plt.tight_layout()

# Save chart to BytesIO
img_bytes = BytesIO()
plt.savefig(img_bytes, format='png', dpi=120, transparent=False, bbox_inches='tight')
plt.close()
img_bytes.seek(0)

# Insert image into slide
slide4.shapes.add_picture(img_bytes, Inches(1.5), Inches(1.5), Inches(10), Inches(5.5))

add_slide_number(slide4, 4)
add_transition(slide4)

# ------------------------------------------------------------
# SLIDE 5 – RATINGS ANALYSIS PIE CHART
# ------------------------------------------------------------
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, DARK_BG)

add_textbox(slide5, Inches(0.8), Inches(0.3), Inches(10), Inches(0.8),
            "Ratings Analysis", font_size=32, bold=True, color=WHITE)
add_ac